"""Extract text content from downloaded Slack files (images, PDFs, Office docs)."""

import asyncio
import json
import logging
import time
from pathlib import Path

log = logging.getLogger(__name__)

IMAGE_TYPES = {"png", "jpg", "jpeg", "gif", "webp", "heic", "bmp", "tiff", "svg"}
MIME_MAP = {
    "png": "image/png", "jpg": "image/jpeg", "jpeg": "image/jpeg",
    "gif": "image/gif", "webp": "image/webp", "heic": "image/heic",
    "bmp": "image/bmp", "tiff": "image/tiff", "svg": "image/svg+xml",
}
TEXT_TYPES = {
    "text", "markdown", "python", "javascript", "typescript", "sql",
    "json", "yaml", "xml", "csv", "shell", "go", "rust", "java",
    "c", "cpp", "html", "css", "plain", "email",
}

VISION_PROMPT = (
    "Describe this image in detail. If it contains text (screenshots, documents, "
    "diagrams, error messages, code), transcribe all visible text. If it's a chart "
    "or graph, describe the data and trends. Focus on information content."
)


def extract_pdf(file_path: Path) -> str | None:
    """Extract text from a PDF file."""
    from pypdf import PdfReader

    try:
        reader = PdfReader(file_path)
        pages = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                pages.append(text.strip())
        content = "\n\n".join(pages)
        return content if len(content) >= 50 else None
    except Exception as e:
        log.debug("PDF extraction failed for %s: %s", file_path.name, e)
        return None


def extract_docx(file_path: Path) -> str | None:
    """Extract text from a DOCX file."""
    from docx import Document

    try:
        doc = Document(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        content = "\n".join(paragraphs)
        return content if len(content) >= 50 else None
    except Exception as e:
        log.debug("DOCX extraction failed for %s: %s", file_path.name, e)
        return None


def extract_xlsx(file_path: Path) -> str | None:
    """Extract text from an XLSX file."""
    from openpyxl import load_workbook

    try:
        wb = load_workbook(file_path, read_only=True, data_only=True)
        lines = []
        for sheet in wb.worksheets:
            lines.append(f"## Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                vals = [str(v) if v is not None else "" for v in row]
                if any(vals):
                    lines.append("\t".join(vals))
        wb.close()
        content = "\n".join(lines)
        return content if len(content) >= 50 else None
    except Exception as e:
        log.debug("XLSX extraction failed for %s: %s", file_path.name, e)
        return None


def extract_pptx(file_path: Path) -> str | None:
    """Extract text from a PPTX file."""
    from pptx import Presentation

    try:
        prs = Presentation(file_path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
            if texts:
                slides.append(f"## Slide {i}\n" + "\n".join(texts))
        content = "\n\n".join(slides)
        return content if len(content) >= 50 else None
    except Exception as e:
        log.debug("PPTX extraction failed for %s: %s", file_path.name, e)
        return None


def extract_text(file_path: Path) -> str | None:
    """Read a text file."""
    try:
        content = file_path.read_text(errors="replace").strip()
        return content if len(content) >= 50 else None
    except Exception as e:
        log.debug("Text read failed for %s: %s", file_path.name, e)
        return None


async def describe_image(
    file_path: Path,
    file_type: str,
    client: "genai.Client",
    model: str,
    sem: asyncio.Semaphore,
) -> str | None:
    """Describe an image using Gemini Vision."""
    from google.genai import types

    mime = MIME_MAP.get(file_type, f"image/{file_type}")
    try:
        data = file_path.read_bytes()
    except Exception as e:
        log.debug("Failed to read image %s: %s", file_path.name, e)
        return None

    try:
        async with sem:
            response = await client.aio.models.generate_content(
                model=model,
                contents=[
                    VISION_PROMPT,
                    types.Part.from_bytes(data=data, mime_type=mime),
                ],
            )
        return response.text if response.text else None
    except Exception as e:
        log.warning("Gemini Vision failed for %s: %s", file_path.name, e)
        return None


async def extract_all(
    files_jsonl: Path,
    raw_dir: Path,
    output_path: Path,
    gemini_api_key: str,
    vision_model: str = "gemini-2.0-flash",
    vision_concurrency: int = 5,
) -> None:
    """Extract content from all downloaded files."""
    from google import genai

    records = []
    with open(files_jsonl) as f:
        for line in f:
            if line.strip():
                records.append(json.loads(line))

    log.info("Loaded %d file records", len(records))

    client = genai.Client(api_key=gemini_api_key)
    sem = asyncio.Semaphore(vision_concurrency)
    start = time.monotonic()

    stats = {"images": 0, "pdfs": 0, "docx": 0, "xlsx": 0, "pptx": 0, "text": 0, "skipped": 0, "failed": 0}
    results: list[dict] = []

    image_records = []
    for rec in records:
        ft = rec.get("file_type", "")
        local_path = raw_dir / rec.get("local_path", "")

        if not local_path.exists():
            stats["skipped"] += 1
            continue

        if ft in IMAGE_TYPES:
            image_records.append(rec)
            continue

        content = None
        method = "text"
        if ft == "pdf":
            content = extract_pdf(local_path)
            method = "pdf"
            stats["pdfs"] += 1
        elif ft == "docx":
            content = extract_docx(local_path)
            method = "docx"
            stats["docx"] += 1
        elif ft in ("xlsx", "xls"):
            content = extract_xlsx(local_path)
            method = "xlsx"
            stats["xlsx"] += 1
        elif ft == "pptx":
            content = extract_pptx(local_path)
            method = "pptx"
            stats["pptx"] += 1
        elif ft in TEXT_TYPES:
            content = extract_text(local_path)
            method = "text"
            stats["text"] += 1
        else:
            stats["skipped"] += 1
            continue

        if content:
            results.append({**rec, "content": content, "extraction_method": method})
        else:
            stats["failed"] += 1

    log.info(
        "Local extraction done: %d PDFs, %d DOCX, %d XLSX, %d PPTX, %d text. %d images pending.",
        stats["pdfs"], stats["docx"], stats["xlsx"], stats["pptx"], stats["text"], len(image_records),
    )

    async def _process_image(rec: dict) -> dict | None:
        ft = rec.get("file_type", "")
        local_path = raw_dir / rec.get("local_path", "")
        content = await describe_image(local_path, ft, client, vision_model, sem)
        if content:
            return {**rec, "content": content, "extraction_method": "vision"}
        return None

    if image_records:
        coros = [_process_image(rec) for rec in image_records]
        done = 0
        for coro in asyncio.as_completed(coros):
            result = await coro
            if result:
                results.append(result)
                stats["images"] += 1
            else:
                stats["failed"] += 1
            done += 1
            if done % 50 == 0:
                elapsed = time.monotonic() - start
                log.info("[%d/%d images] extracted=%d (%.0fs)", done, len(image_records), stats["images"], elapsed)

    with open(output_path, "w") as f:
        for rec in results:
            f.write(json.dumps(rec) + "\n")

    elapsed = time.monotonic() - start
    log.info(
        "Extraction complete: %d records written (%.0fs). "
        "images=%d, pdfs=%d, docx=%d, xlsx=%d, pptx=%d, text=%d, skipped=%d, failed=%d",
        len(results), elapsed,
        stats["images"], stats["pdfs"], stats["docx"], stats["xlsx"],
        stats["pptx"], stats["text"], stats["skipped"], stats["failed"],
    )


def extract(
    files_jsonl: Path,
    raw_dir: Path,
    output_path: Path,
    gemini_api_key: str,
    vision_model: str = "gemini-2.0-flash",
    vision_concurrency: int = 5,
) -> None:
    """Sync entry point."""
    asyncio.run(extract_all(files_jsonl, raw_dir, output_path, gemini_api_key, vision_model, vision_concurrency))
